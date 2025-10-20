#!/usr/bin/env python3
"""
Создание презентации PowerPoint для проекта внедрения AI-ассистента в HR
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np

class HR_AI_Presentation:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_layouts()
        
    def setup_slide_layouts(self):
        """Настройка макетов слайдов"""
        # Используем стандартные макеты PowerPoint
        self.title_layout = self.prs.slide_layouts[0]  # Титульный слайд
        self.content_layout = self.prs.slide_layouts[1]  # Заголовок и содержимое
        self.two_content_layout = self.prs.slide_layouts[3]  # Два содержимых
        self.blank_layout = self.prs.slide_layouts[6]  # Пустой слайд
        
    def add_title_slide(self):
        """Добавление титульного слайда"""
        slide = self.prs.slides.add_slide(self.title_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Внедрение AI-ассистента в HR-процессы"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        
        # Подзаголовок
        subtitle = slide.placeholders[1]
        subtitle.text = "Моделирование снижения расходов и роста продуктивности\n\nАнализ эффективности внедрения Generative AI / LLM (Copilot-стиль)\nв сферах: Канцелярия, офис, бэк-офис"
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
        
    def add_agenda_slide(self):
        """Добавление слайда с программой"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Программа презентации"
        
        # Содержимое
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        agenda_items = [
            "1. Анализ текущего состояния HR-процессов",
            "2. Технологическое решение: AI-ассистент",
            "3. Финансовая модель и ROI",
            "4. Анализ по отделам",
            "5. План внедрения (Диаграмма Ганта)",
            "6. Ожидаемые результаты и риски",
            "7. Рекомендации и следующие шаги"
        ]
        
        for item in agenda_items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_after = Pt(12)
            
    def add_current_state_slide(self):
        """Слайд с анализом текущего состояния"""
        slide = self.prs.slides.add_slide(self.two_content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Текущее состояние HR-процессов"
        
        # Левая колонка - проблемы
        left_content = slide.placeholders[1]
        left_content.text = "Основные проблемы:"
        left_tf = left_content.text_frame
        left_tf.clear()
        
        problems = [
            "• 60% времени тратится на рутинные задачи",
            "• Высокая нагрузка на HR-специалистов",
            "• Медленная обработка документов",
            "• Ошибки в ручном вводе данных",
            "• Недостаток времени на стратегические задачи",
            "• Дублирование работы между отделами"
        ]
        
        for problem in problems:
            p = left_tf.add_paragraph()
            p.text = problem
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            
        # Правая колонка - статистика
        right_content = slide.placeholders[2]
        right_content.text = "Ключевые показатели:"
        right_tf = right_content.text_frame
        right_tf.clear()
        
        stats = [
            "• HR-команда: 8 человек",
            "• Годовые затраты: 960,000 руб",
            "• Время на документооборот: 40%",
            "• Время на коммуникацию: 25%",
            "• Время на аналитику: 20%",
            "• Время на рекрутинг: 15%"
        ]
        
        for stat in stats:
            p = right_tf.add_paragraph()
            p.text = stat
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            
    def add_technology_solution_slide(self):
        """Слайд с технологическим решением"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "AI-ассистент: Технологическое решение"
        
        # Содержимое
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        solution_text = """
        🚀 Generative AI / LLM (Copilot-стиль)
        
        Ключевые возможности:
        • Автоматическое создание HR-документов
        • Интеллектуальная обработка запросов сотрудников
        • Анализ и генерация отчетов
        • Планирование и координация процессов
        • 24/7 поддержка через чат-бот
        
        Области применения:
        • Канцелярия и документооборот
        • Офисные административные задачи
        • Бэк-офис операции
        • HR-аналитика и отчетность
        """
        
        p = tf.add_paragraph()
        p.text = solution_text
        p.font.size = Pt(18)
        
    def add_financial_model_slide(self):
        """Слайд с финансовой моделью"""
        slide = self.prs.slides.add_slide(self.two_content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Финансовая модель и ROI"
        
        # Левая колонка - затраты
        left_content = slide.placeholders[1]
        left_content.text = "Затраты на внедрение:"
        left_tf = left_content.text_frame
        left_tf.clear()
        
        costs = [
            "• Лицензии AI (год): 240,000 руб",
            "• Внедрение: 500,000 руб",
            "• Обучение: 200,000 руб",
            "• Итого в первый год: 940,000 руб",
            "• Ежегодно: 240,000 руб"
        ]
        
        for cost in costs:
            p = left_tf.add_paragraph()
            p.text = cost
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            
        # Правая колонка - экономия
        right_content = slide.placeholders[2]
        right_content.text = "Экономия и выгоды:"
        right_tf = right_content.text_frame
        right_tf.clear()
        
        savings = [
            "• Экономия на рутине: 374,400 руб/год",
            "• Рост продуктивности: 144,000 руб/год",
            "• Общая экономия: 518,400 руб/год",
            "• Чистая экономия: 278,400 руб/год",
            "• ROI за 3 года: 89%",
            "• Окупаемость: 2.5 года"
        ]
        
        for saving in savings:
            p = right_tf.add_paragraph()
            p.text = saving
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            
    def add_department_analysis_slide(self):
        """Слайд с анализом по отделам"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Анализ эффективности по отделам"
        
        # Содержимое
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        dept_analysis = """
        📊 Результаты по отделам:
        
        🎯 HR-аналитика: ROI 156% (лучший результат)
        • Экономия: 37,800 руб/год
        • Улучшение продуктивности: 27%
        
        📋 Кадровый учет: ROI 89%
        • Экономия: 64,000 руб/год
        • Автоматизация: 80% задач
        
        👥 Рекрутинг: ROI 78%
        • Экономия: 78,000 руб/год
        • Ускорение процессов: 60%
        
        💰 Компенсации: ROI 65%
        • Экономия: 41,400 руб/год
        • Точность расчетов: +30%
        
        🎓 Обучение: ROI 45%
        • Экономия: 26,400 руб/год
        • Качество материалов: +40%
        """
        
        p = tf.add_paragraph()
        p.text = dept_analysis
        p.font.size = Pt(16)
        
    def add_implementation_plan_slide(self):
        """Слайд с планом внедрения"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "План внедрения (Диаграмма Ганта)"
        
        # Содержимое
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        plan_text = """
        📅 Временные рамки внедрения: 6 месяцев
        
        Этапы реализации:
        
        1️⃣ Анализ процессов (январь 2024)
        • Аудит текущих HR-процессов
        • Выявление точек автоматизации
        
        2️⃣ Выбор платформы (февраль 2024)
        • Тестирование AI-решений
        • Выбор оптимальной платформы
        
        3️⃣ Техническая подготовка (февраль-март 2024)
        • Настройка инфраструктуры
        • Интеграция с существующими системами
        
        4️⃣ Обучение команды (март 2024)
        • Тренинги по работе с AI
        • Разработка новых процессов
        
        5️⃣ Пилотное внедрение (апрель-май 2024)
        • Тестирование на ограниченной группе
        • Сбор обратной связи
        
        6️⃣ Полное внедрение (июнь 2024)
        • Запуск для всей HR-команды
        • Мониторинг и оптимизация
        """
        
        p = tf.add_paragraph()
        p.text = plan_text
        p.font.size = Pt(16)
        
    def add_expected_results_slide(self):
        """Слайд с ожидаемыми результатами"""
        slide = self.prs.slides.add_slide(self.two_content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Ожидаемые результаты и риски"
        
        # Левая колонка - результаты
        left_content = slide.placeholders[1]
        left_content.text = "Ожидаемые результаты:"
        left_tf = left_content.text_frame
        left_tf.clear()
        
        results = [
            "✅ Сокращение времени на рутину: 65%",
            "✅ Повышение точности: 95%+",
            "✅ Ускорение ответов: 80%",
            "✅ Освобождение времени на стратегию: 40%",
            "✅ Улучшение качества документов",
            "✅ Повышение удовлетворенности сотрудников",
            "✅ Снижение операционных ошибок"
        ]
        
        for result in results:
            p = left_tf.add_paragraph()
            p.text = result
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            
        # Правая колонка - риски
        right_content = slide.placeholders[2]
        right_content.text = "Потенциальные риски:"
        right_tf = right_content.text_frame
        right_tf.clear()
        
        risks = [
            "⚠️ Сопротивление изменениям",
            "⚠️ Необходимость обучения персонала",
            "⚠️ Зависимость от технологий",
            "⚠️ Проблемы с интеграцией",
            "⚠️ Вопросы безопасности данных",
            "⚠️ Необходимость технической поддержки",
            "⚠️ Изменение рабочих процессов"
        ]
        
        for risk in risks:
            p = right_tf.add_paragraph()
            p.text = risk
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            
    def add_recommendations_slide(self):
        """Слайд с рекомендациями"""
        slide = self.prs.slides.add_slide(self.content_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Рекомендации и следующие шаги"
        
        # Содержимое
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        recommendations = """
        🎯 Ключевые рекомендации:
        
        1. Начать с пилотного проекта в отделе HR-аналитики
        • Наибольший потенциал ROI (156%)
        • Относительно простая автоматизация
        
        2. Поэтапное внедрение по отделам
        • Снижение рисков
        • Возможность корректировки подхода
        
        3. Инвестиции в обучение команды
        • Критически важно для успеха
        • Планировать 20% времени на обучение
        
        4. Создание центра компетенций
        • Внутренние эксперты по AI
        • Поддержка и развитие решения
        
        5. Мониторинг и оптимизация
        • Регулярная оценка эффективности
        • Постоянное улучшение процессов
        
        📈 Ожидаемый результат: ROI 89% за 3 года
        """
        
        p = tf.add_paragraph()
        p.text = recommendations
        p.font.size = Pt(16)
        
    def add_conclusion_slide(self):
        """Заключительный слайд"""
        slide = self.prs.slides.add_slide(self.title_layout)
        
        # Заголовок
        title = slide.shapes.title
        title.text = "Заключение"
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)
        
        # Подзаголовок
        subtitle = slide.placeholders[1]
        subtitle.text = """
        Внедрение AI-ассистента в HR-процессы обеспечит:
        
        💰 Экономию 278,400 руб/год
        📈 ROI 89% за 3 года
        ⚡ Повышение продуктивности на 25%
        🎯 Освобождение времени для стратегических задач
        
        Рекомендуется начать с пилотного проекта
        и поэтапно масштабировать решение.
        """
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
        
    def create_presentation(self):
        """Создание полной презентации"""
        print("Создание презентации...")
        
        # Добавление всех слайдов
        self.add_title_slide()
        self.add_agenda_slide()
        self.add_current_state_slide()
        self.add_technology_solution_slide()
        self.add_financial_model_slide()
        self.add_department_analysis_slide()
        self.add_implementation_plan_slide()
        self.add_expected_results_slide()
        self.add_recommendations_slide()
        self.add_conclusion_slide()
        
        # Сохранение презентации
        self.prs.save('/workspace/HR_AI_Implementation_Presentation.pptx')
        print("Презентация сохранена: HR_AI_Implementation_Presentation.pptx")

if __name__ == "__main__":
    # Создание презентации
    presentation = HR_AI_Presentation()
    presentation.create_presentation()
    
    print("\nПрезентация успешно создана!")
    print("Файл: HR_AI_Implementation_Presentation.pptx")