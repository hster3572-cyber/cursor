import os
import math
from datetime import datetime, timedelta
from dataclasses import dataclass
from typing import List, Dict, Tuple

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# -----------------------------
# Configuration and assumptions
# -----------------------------
CURRENCY = "$"
DISCOUNT_RATE = 0.10  # 10% per annum for NPV
WORK_HOURS_PER_FTE = 2080

# Departments and assumptions (Back-office focus)
DEPARTMENTS = [
    {"name": "HR (Кадры)", "fte": 20, "cost_per_fte": 70000, "eligible_share": 0.40, "efficiency": 0.30},
    {"name": "Финансы", "fte": 25, "cost_per_fte": 80000, "eligible_share": 0.30, "efficiency": 0.20},
    {"name": "Юридический/Комплаенс", "fte": 10, "cost_per_fte": 100000, "eligible_share": 0.20, "efficiency": 0.12},
    {"name": "IT Helpdesk", "fte": 15, "cost_per_fte": 75000, "eligible_share": 0.25, "efficiency": 0.15},
    {"name": "Офис/Администрация", "fte": 15, "cost_per_fte": 50000, "eligible_share": 0.50, "efficiency": 0.35},
]

# Adoption ramp and conversion of time saved to direct cost reduction
ADOPTION_RAMP_Q = [0.30, 0.60, 0.80, 0.90]  # operational adoption used for realized savings each quarter
CONVERSION_TO_COST_SAVINGS = 0.60  # share of time savings that converts into payroll reduction

# Seats and pricing (based on public pricing circa 2024-2025)
SEATS_TOTAL = 120
LICENSE_PRICE_PER_USER_PER_MONTH = 30  # e.g., Microsoft Copilot for M365
API_PRICE_PER_USER_PER_MONTH = 10      # average LLM API ops and orchestration
SUPPORT_COST_PER_YEAR = 60000          # prompt eng./ops/monitoring
IMPLEMENTATION_SERVICES = 100000       # one-time
TRAINING_PER_SEAT = 1000               # one-time

# Seat ramp for Year 1 (procurement pace)
SEAT_RAMP_Q = [0.50, 0.75, 1.00, 1.00]  # share of SEATS_TOTAL licensed per quarter

# Year 2 and 3 steady-state adoptions for savings
ADOPTION_STEADY_Y2 = 0.90
ADOPTION_STEADY_Y3 = 0.95

# Gantt plan (weeks), starting next month 1st day
TODAY = datetime.today()
START_DATE = (TODAY.replace(day=1) + timedelta(days=32)).replace(day=1)  # first day of next month
GANTT_TASKS = [
    ("Дисквери и приоритизация кейсов", 0, 3),
    ("Доступ к данным и безопасность", 1, 4),
    ("Пилот: HR+Администрация", 3, 5),
    ("UAT и обучение: волна 1", 6, 3),
    ("Масштабирование: HR/Адм", 8, 8),
    ("Расширение: Финансы и IT", 12, 12),
    ("Юрид/Комплаенс интеграция", 16, 8),
    ("Оптимизация и аналитика", 20, 6),
]


# -----------------------------
# Helpers
# -----------------------------

def fmt_currency(v: float) -> str:
    return f"{CURRENCY}{v:,.0f}" if not math.isnan(v) else "-"


def fmt_pct(p: float) -> str:
    return f"{p*100:.0f}%"


@dataclass
class DeptResult:
    name: str
    fte: int
    cost_per_fte: float
    baseline_cost: float
    eligible_share: float
    efficiency: float
    y1_cost_savings: float
    y2_cost_savings: float
    y3_cost_savings: float
    y1_fte_freed: float
    y2_fte_freed: float
    y3_fte_freed: float


# -----------------------------
# Core calculations
# -----------------------------

def compute_department_results() -> Tuple[pd.DataFrame, Dict[str, float]]:
    results: List[DeptResult] = []

    for d in DEPARTMENTS:
        baseline_cost = d["fte"] * d["cost_per_fte"]
        # Year 1: quarter-by-quarter
        y1_savings = 0.0
        y1_fte_freed = 0.0
        for a in ADOPTION_RAMP_Q:
            q_savings = baseline_cost * d["eligible_share"] * d["efficiency"] * a * CONVERSION_TO_COST_SAVINGS * 0.25
            y1_savings += q_savings
            q_fte = d["fte"] * d["eligible_share"] * d["efficiency"] * a * 0.25
            y1_fte_freed += q_fte

        # Year 2 and 3 steady-state
        y2_savings = baseline_cost * d["eligible_share"] * d["efficiency"] * ADOPTION_STEADY_Y2 * CONVERSION_TO_COST_SAVINGS
        y3_savings = baseline_cost * d["eligible_share"] * d["efficiency"] * ADOPTION_STEADY_Y3 * CONVERSION_TO_COST_SAVINGS

        y2_fte_freed = d["fte"] * d["eligible_share"] * d["efficiency"] * ADOPTION_STEADY_Y2
        y3_fte_freed = d["fte"] * d["eligible_share"] * d["efficiency"] * ADOPTION_STEADY_Y3

        results.append(
            DeptResult(
                name=d["name"],
                fte=d["fte"],
                cost_per_fte=d["cost_per_fte"],
                baseline_cost=baseline_cost,
                eligible_share=d["eligible_share"],
                efficiency=d["efficiency"],
                y1_cost_savings=y1_savings,
                y2_cost_savings=y2_savings,
                y3_cost_savings=y3_savings,
                y1_fte_freed=y1_fte_freed,
                y2_fte_freed=y2_fte_freed,
                y3_fte_freed=y3_fte_freed,
            )
        )

    df = pd.DataFrame([r.__dict__ for r in results])

    # Costs
    y1_license = 0.0
    y1_api = 0.0
    for seat_share in SEAT_RAMP_Q:
        q_seats = SEATS_TOTAL * seat_share
        y1_license += q_seats * LICENSE_PRICE_PER_USER_PER_MONTH * 3  # 3 months per quarter
        y1_api += q_seats * API_PRICE_PER_USER_PER_MONTH * 3

    y1_support = SUPPORT_COST_PER_YEAR
    y1_one_time = IMPLEMENTATION_SERVICES + SEATS_TOTAL * TRAINING_PER_SEAT

    y2_license = SEATS_TOTAL * LICENSE_PRICE_PER_USER_PER_MONTH * 12
    y2_api = SEATS_TOTAL * API_PRICE_PER_USER_PER_MONTH * 12
    y2_support = SUPPORT_COST_PER_YEAR

    # Assume similar recurring in Y3
    y3_license = y2_license
    y3_api = y2_api
    y3_support = y2_support

    totals = {
        "y1_savings_total": float(df["y1_cost_savings"].sum()),
        "y2_savings_total": float(df["y2_cost_savings"].sum()),
        "y3_savings_total": float(df["y3_cost_savings"].sum()),
        "y1_recurring_costs": y1_license + y1_api + y1_support,
        "y1_one_time_costs": y1_one_time,
        "y2_recurring_costs": y2_license + y2_api + y2_support,
        "y3_recurring_costs": y3_license + y3_api + y3_support,
    }

    # Net cash flows
    y1_net = totals["y1_savings_total"] - (totals["y1_recurring_costs"] + totals["y1_one_time_costs"])
    y2_net = totals["y2_savings_total"] - totals["y2_recurring_costs"]
    y3_net = totals["y3_savings_total"] - totals["y3_recurring_costs"]

    totals.update({"y1_net": y1_net, "y2_net": y2_net, "y3_net": y3_net})

    # Payback (rough, annual buckets)
    cumulative = y1_net + y2_net
    if y1_net >= 0:
        payback_years = 0.5  # within Y1; rough placeholder when positive
    elif y1_net < 0 and cumulative >= 0:
        frac = -y1_net / (y2_net if y2_net != 0 else 1)
        payback_years = 1 + frac
    else:
        # would occur in Y3
        denom = (y2_net + y3_net) if (y2_net + y3_net) != 0 else 1
        frac = -(y1_net + y2_net) / (y3_net if y3_net != 0 else 1)
        payback_years = 2 + max(0.0, frac)

    totals["payback_years"] = float(max(0.0, payback_years))

    # NPV over 3 years
    npv = (y1_net / ((1 + DISCOUNT_RATE) ** 1)) + (y2_net / ((1 + DISCOUNT_RATE) ** 2)) + (y3_net / ((1 + DISCOUNT_RATE) ** 3))
    totals["npv_3y"] = float(npv)

    return df, totals


# -----------------------------
# Gantt chart
# -----------------------------

def build_gantt_chart(path: str) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)

    fig, ax = plt.subplots(figsize=(10, 4.5))

    task_names = [t[0] for t in GANTT_TASKS]
    starts = [START_DATE + timedelta(weeks=t[1]) for t in GANTT_TASKS]
    ends = [START_DATE + timedelta(weeks=t[1] + t[2]) for t in GANTT_TASKS]

    # Plot bars (reverse to show first on top)
    y_pos = np.arange(len(task_names))
    for i, (name, s, e) in enumerate(zip(task_names, starts, ends)):
        ax.barh(
            y_pos[i],
            mdates.date2num(e) - mdates.date2num(s),
            left=mdates.date2num(s),
            height=0.5,
            color="#4F81BD",
            edgecolor="#274B7A",
        )

    ax.set_yticks(y_pos)
    ax.set_yticklabels(task_names)

    ax.xaxis_date()
    ax.xaxis.set_major_locator(mdates.MonthLocator(interval=1))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%b '%y"))
    ax.grid(axis="x", linestyle="--", alpha=0.4)
    ax.set_title("Диаграмма Ганта: внедрение AI-ассистента")
    fig.tight_layout()
    plt.savefig(path, dpi=200)
    plt.close(fig)


# -----------------------------
# PPT building helpers
# -----------------------------

def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullets_slide(prs: Presentation, title: str, bullets: List[str]):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = b
        p.level = 0


def add_table_slide(prs: Presentation, title: str, table_data: List[List[str]]):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    rows = len(table_data)
    cols = len(table_data[0]) if rows else 0
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.3 * rows)

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = tbl_shape.table

    # Header formatting
    for c in range(cols):
        cell = table.cell(0, c)
        cell.text = table_data[0][c]
        for p in cell.text_frame.paragraphs:
            p.font.bold = True

    # Body
    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = str(table_data[r][c])



def add_image_slide(prs: Presentation, title: str, image_path: str):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.2)
    slide.shapes.add_picture(image_path, left, top, height=Inches(5))


# -----------------------------
# Build presentation
# -----------------------------

def build_presentation(output_path: str):
    df, totals = compute_department_results()

    prs = Presentation()

    # 1) Title
    add_title_slide(
        prs,
        "AI Copilot для Back-Office и HR",
        "Модель снижения расходов и роста продуктивности (до/после по отделам)",
    )

    # 2) Executive summary
    bullets = [
        f"Год 1: экономия {fmt_currency(totals['y1_savings_total'])}, издержки {fmt_currency(totals['y1_recurring_costs'] + totals['y1_one_time_costs'])}, итог {fmt_currency(totals['y1_net'])}",
        f"Год 2: экономия {fmt_currency(totals['y2_savings_total'])}, издержки {fmt_currency(totals['y2_recurring_costs'])}, итог {fmt_currency(totals['y2_net'])}",
        f"Год 3: экономия {fmt_currency(totals['y3_savings_total'])}, издержки {fmt_currency(totals['y3_recurring_costs'])}, итог {fmt_currency(totals['y3_net'])}",
        f"Окупаемость ~ {totals['payback_years']:.1f} года; NPV (3 года, 10%) = {fmt_currency(totals['npv_3y'])}",
    ]
    add_bullets_slide(prs, "Краткая сводка", bullets)

    # 3) Use-cases per function
    use_cases = [
        "HR: авто-драфт офферов, писем, онбординг-документов; Q&A по политикам",
        "Финансы: сбор данных для отчетов, сверка транзакций, черновики писем поставщикам",
        "Юрид/Комплаенс: генерация черновиков договоров/допсоглашений, поиск рисков",
        "IT Helpdesk: ответы на типовые тикеты, базы знания, инструкции",
        "Администрация/Офис: календарь/встречи, протоколы, заявки, запросы к поставщикам",
    ]
    add_bullets_slide(prs, "Приоритетные кейсы автоматизации (Copilot-стиль)", use_cases)

    # 4) Assumptions table
    assumptions = [
        ["Параметр", "Значение"],
        ["Конверсия времени в снижение ФОТ", fmt_pct(CONVERSION_TO_COST_SAVINGS)],
        ["Лицензии (в мес.)", fmt_currency(LICENSE_PRICE_PER_USER_PER_MONTH)],
        ["API (в мес.)", fmt_currency(API_PRICE_PER_USER_PER_MONTH)],
        ["Сервис поддержки (в год)", fmt_currency(SUPPORT_COST_PER_YEAR)],
        ["Внедрение (one-time)", fmt_currency(IMPLEMENTATION_SERVICES)],
        ["Обучение на сотрудника (one-time)", fmt_currency(TRAINING_PER_SEAT)],
        ["Лицензий (итог)", f"{SEATS_TOTAL}"],
        ["Принятый дисконт на NPV", f"{int(DISCOUNT_RATE*100)}%"],
    ]
    add_table_slide(prs, "Ключевые допущения и стоимость", assumptions)

    # 5) Before/After by department (Year 1)
    table = [[
        "Отдел",
        "FTE",
        "ФОТ (база)",
        "Доля рутины",
        "Эффект ИИ",
        "Экономия Г1",
        "Освобождено FTE Г1",
    ]]
    for _, row in df.iterrows():
        table.append([
            row["name"],
            f"{int(row['fte'])}",
            fmt_currency(row["baseline_cost"]),
            fmt_pct(row["eligible_share"]),
            fmt_pct(row["efficiency"]),
            fmt_currency(row["y1_cost_savings"]),
            f"{row['y1_fte_freed']:.2f}",
        ])
    add_table_slide(prs, "До/после по отделам — Год 1", table)

    # 6) Financial summary (Y1-Y3)
    fin_table = [
        ["Показатель", "Год 1", "Год 2", "Год 3"],
        ["Экономия (всего)", fmt_currency(totals['y1_savings_total']), fmt_currency(totals['y2_savings_total']), fmt_currency(totals['y3_savings_total'])],
        ["Издержки (повторяющиеся)", fmt_currency(totals['y1_recurring_costs']), fmt_currency(totals['y2_recurring_costs']), fmt_currency(totals['y3_recurring_costs'])],
        ["Издержки (разовые)", fmt_currency(totals['y1_one_time_costs']), fmt_currency(0), fmt_currency(0)],
        ["Итоговый денежный поток", fmt_currency(totals['y1_net']), fmt_currency(totals['y2_net']), fmt_currency(totals['y3_net'])],
    ]
    add_table_slide(prs, "Сводка финансовых результатов", fin_table)

    # 7) FTE capacity uplift (Y1)
    cap_table = [["Отдел", "Свободная емкость (FTE экв., Г1)"]]
    for _, row in df.iterrows():
        cap_table.append([row["name"], f"{row['y1_fte_freed']:.2f}"])
    add_table_slide(prs, "Рост продуктивности: эквивалент FTE (Год 1)", cap_table)

    # 8) Gantt chart
    gantt_path = os.path.join("assets", "gantt.png")
    build_gantt_chart(gantt_path)
    add_image_slide(prs, "План внедрения (Гант)", gantt_path)

    # 9) Risks/Mitigations
    risks = [
        "Качество данных и доступов — выстраивание безопасных коннекторов и аудит логов",
        "Принятие пользователями — обучение, FAQ, внутренние чемпионы",
        "Точность LLM — тест-кейсы, guardrails, human-in-the-loop",
        "Соблюдение политики ИБ/Юрид — встраивание DLP, анонимизация, ограничение PII",
    ]
    add_bullets_slide(prs, "Основные риски и меры", risks)

    # 10) References slide
    refs = [
        "Pricing: Microsoft Copilot for M365 — $30/польз./мес (2024)",
        "ChatGPT Team/Enterprise, Claude Team — ориентиры $25–$60/польз./мес",
        "McKinsey (2023): потенциал GenAI — 20–30% высвобождение времени для бэк-офиса",
        "Deloitte, Gartner (2023–2024): пилоты 6–12 недель, масштабирование 2–3 месяца",
    ]
    add_bullets_slide(prs, "Источники и ориентиры рынка", refs)

    # Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True) if os.path.dirname(output_path) else None
    prs.save(output_path)


if __name__ == "__main__":
    out_pptx = os.path.join("HR_Automation_AI_Impact.pptx")
    build_presentation(out_pptx)
    print(f"Presentation generated: {out_pptx}")
