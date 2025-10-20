#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Генератор презентации для анализа внедрения AI-ассистента в HR
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
from io import BytesIO
import base64

class PresentationGenerator:
    """Класс для создания презентации по результатам анализа"""
    
    def __init__(self):
        self.prs = Presentation()
        self.setup_theme()
    
    def setup_theme(self):
        """Настройка темы презентации"""
        # Установка размера слайда
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
    
    def add_title_slide(self):
        """Добавление титульного слайда"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Внедрение AI-ассистента для автоматизации административных HR-задач"
        subtitle.text = ("Анализ снижения расходов и роста продуктивности\n"
                        "Моделирование эффектов по отделам\n\n"
                        f"Дата: {pd.Timestamp.now().strftime('%d.%m.%Y')}")
        
        # Форматирование заголовка
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Форматирование подзаголовка
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(64, 64, 64)
    
    def add_agenda_slide(self):
        """Добавление слайда с повесткой дня"""
        slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Повестка дня"
        
        agenda_items = [
            "1. Цели и задачи проекта",
            "2. Анализ текущего состояния",
            "3. Технологическое решение",
            "4. Финансовая модель",
            "5. Анализ воздействия по отделам",
            "6. План внедрения (диаграмма Ганта)",
            "7. Ожидаемые результаты",
            "8. Риски и митигация",
            "9. Выводы и рекомендации"
        ]
        
        content.text = "\n".join(agenda_items)
        
        # Форматирование
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.space_after = Pt(8)
    
    def add_objectives_slide(self):
        """Добавление слайда с целями проекта"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Цели и задачи проекта"
        
        objectives_text = """
ГЛАВНАЯ ЦЕЛЬ:
Автоматизация административных задач с помощью AI-ассистента (Copilot-стиль)

КЛЮЧЕВЫЕ ЗАДАЧИ:
• Снижение операционных расходов на персонал
• Повышение продуктивности сотрудников
• Улучшение качества выполнения рутинных задач
• Сокращение количества ошибок
• Освобождение времени для стратегических задач

ОБЛАСТИ ПРИМЕНЕНИЯ:
• Канцелярия - документооборот, архивирование
• HR отдел - обработка заявлений, отчетность
• Бэк-офис - административные процессы, аналитика
        """
        
        content.text = objectives_text.strip()
        
        # Форматирование
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(6)
    
    def add_current_state_slide(self, dept_comparison_df):
        """Добавление слайда с анализом текущего состояния"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Анализ текущего состояния"
        
        # Создание таблицы
        rows = len(dept_comparison_df) + 1
        cols = 4
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(14)
        height = Inches(5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Заголовки таблицы
        headers = ['Отдел', 'Сотрудники', 'Расходы на админ задачи (руб/мес)', 'Потенциал автоматизации']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Заполнение данных
        for i, row in dept_comparison_df.iterrows():
            table.cell(i+1, 0).text = row['Отдел']
            table.cell(i+1, 1).text = str(row['Сотрудники'])
            table.cell(i+1, 2).text = f"{row['Текущие расходы на админ задачи (руб/мес)']:,.0f}"
            table.cell(i+1, 3).text = f"{row['Потенциал автоматизации (%)']:.0f}%"
            
            # Форматирование ячеек данных
            for j in range(4):
                table.cell(i+1, j).text_frame.paragraphs[0].font.size = Pt(11)
    
    def add_technology_slide(self):
        """Добавление слайда с описанием технологического решения"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Технологическое решение: AI-ассистент"
        
        tech_content = """
ТЕХНОЛОГИЯ: Generative AI / Large Language Model (LLM)

КЛЮЧЕВЫЕ ВОЗМОЖНОСТИ:
• Автоматическая обработка документов
• Генерация отчетов и справок
• Интеллектуальная классификация заявлений
• Автоматизация рутинных запросов
• Помощь в принятии решений на основе данных

ПРЕИМУЩЕСТВА COPILOT-ПОДХОДА:
• Естественное языковое взаимодействие
• Контекстное понимание задач
• Обучение на корпоративных данных
• Интеграция с существующими системами
• Масштабируемость решения

ОЖИДАЕМЫЕ ЭФФЕКТЫ:
• Сокращение времени на рутинные задачи на 75%
• Снижение ошибок на 85%
• Увеличение скорости обработки в 2.5 раза
        """
        
        content.text = tech_content.strip()
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(6)
    
    def add_financial_model_slide(self, financial_df):
        """Добавление слайда с финансовой моделью"""
        slide_layout = self.prs.slide_layouts[5]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Заголовок
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
        title_frame = title_shape.text_frame
        title_frame.text = "Финансовая модель проекта"
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Ключевые показатели
        kpi_data = [
            ("Первоначальные инвестиции", "700,000 руб"),
            ("Ежемесячные расходы", "40,000 руб"),
            ("Ежемесячная экономия", f"{financial_df['monthly_savings'].iloc[0]:,.0f} руб"),
            ("Точка безубыточности", f"{financial_df[financial_df['net_benefit'] >= 0]['month'].iloc[0]} месяц"),
            ("ROI через 12 месяцев", f"{financial_df[financial_df['month'] == 12]['roi_percent'].iloc[0]:.1f}%"),
            ("ROI через 36 месяцев", f"{financial_df[financial_df['month'] == 36]['roi_percent'].iloc[0]:.1f}%")
        ]
        
        # Создание таблицы KPI
        left = Inches(2)
        top = Inches(2)
        width = Inches(12)
        height = Inches(4)
        
        table = slide.shapes.add_table(len(kpi_data), 2, left, top, width, height).table
        
        for i, (metric, value) in enumerate(kpi_data):
            table.cell(i, 0).text = metric
            table.cell(i, 1).text = value
            
            # Форматирование
            table.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(16)
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(16)
            table.cell(i, 1).text_frame.paragraphs[0].font.bold = True
            table.cell(i, 1).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
            
            # Цветовое кодирование
            if "ROI" in metric:
                table.cell(i, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 128, 0)
    
    def add_department_impact_slide(self, dept_comparison_df):
        """Добавление слайда с анализом воздействия по отделам"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Анализ воздействия по отделам"
        
        # Создание расширенной таблицы
        rows = len(dept_comparison_df) + 1
        cols = 5
        
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(15)
        height = Inches(5.5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Заголовки
        headers = ['Отдел', 'Экономия от автоматизации', 'Прирост продуктивности', 
                  'Снижение ошибок', 'Общая выгода (руб/мес)']
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(79, 129, 189)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Заполнение данных
        for i, row in dept_comparison_df.iterrows():
            table.cell(i+1, 0).text = row['Отдел']
            table.cell(i+1, 1).text = f"{row['Экономия от автоматизации (руб/мес)']:,.0f}"
            table.cell(i+1, 2).text = f"{row['Прирост продуктивности (руб/мес)']:,.0f}"
            table.cell(i+1, 3).text = f"{row['Экономия от снижения ошибок (руб/мес)']:,.0f}"
            table.cell(i+1, 4).text = f"{row['Общая выгода (руб/мес)']:,.0f}"
            
            for j in range(5):
                table.cell(i+1, j).text_frame.paragraphs[0].font.size = Pt(10)
                if j > 0:  # Числовые столбцы
                    table.cell(i+1, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def add_implementation_plan_slide(self):
        """Добавление слайда с планом внедрения"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "План внедрения (основные этапы)"
        
        plan_content = """
ФАЗА 1: ПОДГОТОВКА (1-2 месяц)
• Анализ и планирование
• Выбор AI-платформы
• Закупка лицензий

ФАЗА 2: ТЕХНИЧЕСКАЯ РЕАЛИЗАЦИЯ (2-4 месяц)
• Настройка инфраструктуры
• Интеграция с существующими системами
• Разработка процедур и регламентов

ФАЗА 3: ОБУЧЕНИЕ И ПИЛОТ (3-5 месяц)
• Обучение персонала всех отделов
• Пилотное внедрение по отделам
• Анализ результатов пилота

ФАЗА 4: ПОЛНОЕ ВНЕДРЕНИЕ (5-6 месяц)
• Корректировка и оптимизация
• Развертывание на всех отделах
• Мониторинг и поддержка

ОБЩАЯ ПРОДОЛЖИТЕЛЬНОСТЬ: 9 месяцев
        """
        
        content.text = plan_content.strip()
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(8)
    
    def add_expected_results_slide(self):
        """Добавление слайда с ожидаемыми результатами"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Ожидаемые результаты"
        
        results_content = """
КОЛИЧЕСТВЕННЫЕ РЕЗУЛЬТАТЫ:
• Экономия на зарплатном фонде: 1,200,000+ руб/год
• Сокращение времени на админ задачи: 75%
• Снижение количества ошибок: 85%
• Увеличение скорости обработки: в 2.5 раза
• ROI через 12 месяцев: 180%+

КАЧЕСТВЕННЫЕ РЕЗУЛЬТАТЫ:
• Повышение удовлетворенности сотрудников
• Улучшение качества обслуживания
• Стандартизация процессов
• Освобождение времени для стратегических задач
• Повышение конкурентоспособности

ДОЛГОСРОЧНЫЕ ЭФФЕКТЫ:
• Создание базы для дальнейшей цифровизации
• Развитие компетенций в области AI
• Масштабирование решения на другие процессы
        """
        
        content.text = results_content.strip()
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(6)
    
    def add_risks_slide(self):
        """Добавление слайда с рисками и митигацией"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Риски и меры митигации"
        
        risks_content = """
ТЕХНИЧЕСКИЕ РИСКИ:
• Сложность интеграции → Поэтапное внедрение, тестирование
• Качество AI-ответов → Обучение на корпоративных данных
• Безопасность данных → Строгие протоколы безопасности

ОРГАНИЗАЦИОННЫЕ РИСКИ:
• Сопротивление персонала → Программа обучения и поддержки
• Недостаток компетенций → Привлечение экспертов, тренинги
• Изменение процессов → Постепенная адаптация, обратная связь

ФИНАНСОВЫЕ РИСКИ:
• Превышение бюджета → Детальное планирование, контроль
• Задержка ROI → Фокус на быстрых победах
• Скрытые расходы → Резерв 20% от бюджета

ОБЩАЯ ОЦЕНКА РИСКОВ: СРЕДНИЙ УРОВЕНЬ
Все риски управляемы при правильном планировании
        """
        
        content.text = risks_content.strip()
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(15)
            paragraph.space_after = Pt(6)
    
    def add_conclusions_slide(self):
        """Добавление слайда с выводами"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Выводы и рекомендации"
        
        conclusions_content = """
КЛЮЧЕВЫЕ ВЫВОДЫ:
✓ Проект экономически обоснован (ROI 180%+ через год)
✓ Высокий потенциал автоматизации во всех отделах
✓ Быстрая окупаемость (точка безубыточности - 8 месяц)
✓ Значительное повышение эффективности процессов

РЕКОМЕНДАЦИИ:
1. НЕМЕДЛЕННО начать проект - конкурентное преимущество
2. Начать с канцелярии (наибольший потенциал экономии)
3. Инвестировать в обучение персонала
4. Создать центр компетенций по AI
5. Планировать масштабирование на другие процессы

СЛЕДУЮЩИЕ ШАГИ:
• Утверждение проекта и бюджета
• Формирование проектной команды
• Начало фазы анализа и планирования
• Выбор поставщика AI-решения

ПРОЕКТ ГОТОВ К РЕАЛИЗАЦИИ!
        """
        
        content.text = conclusions_content.strip()
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
            paragraph.space_after = Pt(6)
    
    def save_presentation(self, filename="AI_HR_Automation_Analysis.pptx"):
        """Сохранение презентации"""
        self.prs.save(filename)
        return filename

def create_full_presentation(dept_comparison_df, financial_df):
    """Создание полной презентации"""
    generator = PresentationGenerator()
    
    # Добавление всех слайдов
    generator.add_title_slide()
    generator.add_agenda_slide()
    generator.add_objectives_slide()
    generator.add_current_state_slide(dept_comparison_df)
    generator.add_technology_slide()
    generator.add_financial_model_slide(financial_df)
    generator.add_department_impact_slide(dept_comparison_df)
    generator.add_implementation_plan_slide()
    generator.add_expected_results_slide()
    generator.add_risks_slide()
    generator.add_conclusions_slide()
    
    # Сохранение
    filename = generator.save_presentation()
    return filename

if __name__ == "__main__":
    # Пример использования
    from ai_hr_automation_analysis import AIHRAutomationAnalyzer
    
    analyzer = AIHRAutomationAnalyzer()
    financial_df, _, _ = analyzer.create_financial_model()
    dept_comparison = analyzer.create_department_comparison()
    
    filename = create_full_presentation(dept_comparison, financial_df)
    print(f"Презентация создана: {filename}")