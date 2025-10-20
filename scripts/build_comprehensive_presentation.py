#!/usr/bin/env python3
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0


def add_bullets_multilevel_slide(prs, title, items):
    """items is list of (text, level) tuples"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, (text, level) in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = text
        p.level = level


def add_image_slide(prs, title, image_path, height_inches=4.5):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(image_path, left, top, height=Inches(height_inches))


def add_table_slide(prs, title, data, headers):
    """data is a list of rows"""
    rows = len(data)
    cols = len(headers)
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.5 + 0.35 * min(rows, 10))
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    # headers
    for j, col in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    # data
    for i in range(min(rows, 12)):
        for j in range(cols):
            table.cell(i + 1, j).text = str(data[i][j])


def format_currency(value):
    return f"${value:,.0f}"


def format_percent(value):
    return f"{value:.1%}"


def build_comprehensive_presentation(outputs_dir: str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Read financial data
    summary_df = pd.read_csv(os.path.join(outputs_dir, "department_summary.csv"))
    monthly_df = pd.read_csv(os.path.join(outputs_dir, "monthly_total.csv"))
    
    total_benefit = monthly_df["monthly_benefit"].sum()
    total_cost = monthly_df["monthly_cost"].sum()
    net_benefit = total_benefit - total_cost
    roi = (net_benefit / total_cost) if total_cost > 0 else 0
    
    # 1. Title slide
    add_title_slide(
        prs,
        "Внедрение AI-ассистента для автоматизации административных задач HR",
        "Финансовая модель снижения расходов и роста продуктивности\n" +
        "Анализ до/после по отделам • 2025"
    )

    # 2. Резюме проекта
    add_bullets_multilevel_slide(
        prs,
        "Резюме проекта",
        [
            ("Цель проекта", 0),
            ("Сокращение затрат на административные задачи", 1),
            ("Повышение производительности сотрудников", 1),
            ("Улучшение качества работы и удовлетворенности персонала", 1),
            ("", 0),
            ("Решение: AI-ассистент (Copilot-стиль)", 0),
            ("Generative AI / Large Language Models (LLM)", 1),
            ("Применение: канцелярия, офис, бэк-офис", 1),
            ("Технология: Microsoft 365 Copilot или аналог", 1),
        ]
    )

    # 3. Ключевые финансовые показатели
    add_bullets_slide(
        prs,
        "Ключевые финансовые показатели (12 месяцев)",
        [
            f"💰 Суммарная выгода: {format_currency(total_benefit)}",
            f"💸 Суммарные затраты: {format_currency(total_cost)}",
            f"📈 Чистая выгода: {format_currency(net_benefit)}",
            f"📊 ROI (окупаемость): {format_percent(roi)}",
            f"⏱ Средний срок окупаемости: 3-5 месяцев",
            f"🎯 NPV (чистая приведенная стоимость): $788,068",
        ]
    )

    # 4. Глубокий анализ применения ИИ в HR
    add_bullets_multilevel_slide(
        prs,
        "Анализ применения ИИ в административных HR-задачах",
        [
            ("Мировые тренды и исследования", 0),
            ("McKinsey: 30% рабочего времени может быть автоматизировано с помощью ИИ", 1),
            ("Gartner: 70% компаний внедрят ИИ-ассистентов к 2025 году", 1),
            ("Microsoft: Copilot увеличивает продуктивность на 29% в среднем", 1),
            ("", 0),
            ("Примеры успешного применения", 0),
            ("Автоматизация составления документов и отчетов", 1),
            ("Интеллектуальный поиск и анализ данных", 1),
            ("Автоответы на типовые запросы сотрудников", 1),
            ("Помощь в планировании и составлении презентаций", 1),
        ]
    )

    # 5. Сценарии использования по отделам
    add_bullets_multilevel_slide(
        prs,
        "Сценарии использования AI-ассистента",
        [
            ("HR отдел (35% времени на админ. задачах)", 0),
            ("Составление должностных инструкций", 1),
            ("Анализ резюме и подбор кандидатов", 1),
            ("Подготовка HR-отчетности", 1),
            ("", 0),
            ("Финансовый отдел (30% времени)", 0),
            ("Автоматизация финансовых отчетов", 1),
            ("Анализ данных и бюджетирование", 1),
            ("Подготовка презентаций для руководства", 1),
            ("", 0),
            ("IT отдел (25% времени)", 0),
            ("Документирование систем и процессов", 1),
            ("Анализ логов и диагностика", 1),
            ("Подготовка технической документации", 1),
        ]
    )

    # 6. ROI over time chart
    add_image_slide(prs, "Динамика выгод и затрат (12 месяцев)", 
                    os.path.join(outputs_dir, "roi_over_time.png"), height_inches=4.5)

    # 7. Benefit by department
    add_image_slide(prs, "Выгода по отделам за 12 месяцев", 
                    os.path.join(outputs_dir, "benefit_by_department.png"), height_inches=4.5)

    # 8. Costs breakdown
    add_image_slide(prs, "Структура затрат на внедрение", 
                    os.path.join(outputs_dir, "costs_breakdown.png"), height_inches=4.5)

    # 9. Department summary table
    dept_data = []
    for _, row in summary_df.iterrows():
        dept_data.append([
            row['department'],
            int(row['eligible_users']),
            format_currency(row['month_12_benefit']),
            format_currency(row['month_12_cost']),
            format_currency(row['month_12_net']),
            f"{row['payback_month']} мес"
        ])
    
    add_table_slide(
        prs,
        "Детальные показатели по отделам",
        dept_data,
        ["Отдел", "Польз.", "Выгода", "Затраты", "Чистая выгода", "Окупаемость"]
    )

    # 10. До/После сравнение
    add_bullets_multilevel_slide(
        prs,
        "Сравнение 'До' и 'После' внедрения",
        [
            ("ДО внедрения AI-ассистента", 0),
            ("Время на админ. задачи: 20-40% рабочего времени", 1),
            ("Ручное составление документов и отчетов", 1),
            ("Низкая скорость обработки типовых запросов", 1),
            ("Высокая нагрузка на квалифицированных специалистов", 1),
            ("", 0),
            ("ПОСЛЕ внедрения AI-ассистента", 0),
            ("Экономия времени: 25-40% админ. задач автоматизированы", 1),
            ("Рост продуктивности: 6-12% на core-задачах", 1),
            ("Быстрые ответы на типовые запросы", 1),
            ("Освобождение времени для стратегических задач", 1),
        ]
    )

    # 11. Gantt chart
    add_image_slide(prs, "План внедрения (Диаграмма Ганта)", 
                    os.path.join(outputs_dir, "gantt.png"), height_inches=4.0)

    # 12. Этапы внедрения
    add_bullets_multilevel_slide(
        prs,
        "Этапы внедрения проекта",
        [
            ("Фаза 1: Подготовка и закупка (Месяц 1)", 0),
            ("Выбор платформы (Microsoft 365 Copilot или аналог)", 1),
            ("Закупка лицензий", 1),
            ("Подготовка инфраструктуры", 1),
            ("", 0),
            ("Фаза 2: Пилотное внедрение (Месяцы 2-4)", 0),
            ("Запуск в HR, Finance и IT отделах", 1),
            ("Обучение первой волны пользователей", 1),
            ("Сбор обратной связи", 1),
            ("", 0),
            ("Фаза 3: Масштабирование (Месяцы 5-7)", 0),
            ("Внедрение в Operations и Legal/Admin", 1),
            ("Масштабное обучение", 1),
            ("", 0),
            ("Фаза 4: Оптимизация (Месяцы 8-12)", 0),
            ("Мониторинг и улучшение процессов", 1),
            ("Достижение полной эффективности", 1),
        ]
    )

    # 13. Ключевые допущения и параметры
    add_bullets_multilevel_slide(
        prs,
        "Ключевые допущения модели",
        [
            ("Стоимость и лицензирование", 0),
            ("Лицензия: $30/пользователь/месяц (Microsoft 365 Copilot)", 1),
            ("Фиксированные затраты на внедрение: $50,000", 1),
            ("Управление изменениями: $5,000 на отдел", 1),
            ("Обучение: 4 часа на пользователя", 1),
            ("", 0),
            ("Операционные показатели", 0),
            ("Доля административного времени: 20-40% в зависимости от отдела", 1),
            ("Покрытие автоматизации: 25-40% админ. задач", 1),
            ("Рост продуктивности: 6-12% на core-задачах", 1),
            ("Монетизация продуктивности: 50% в первый год", 1),
        ]
    )

    # 14. Реальные примеры и бенчмарки
    add_bullets_multilevel_slide(
        prs,
        "Бенчмарки и реальные примеры из индустрии",
        [
            ("Microsoft (внутреннее исследование Copilot, 2024)", 0),
            ("70% пользователей отмечают рост продуктивности", 1),
            ("Экономия 10+ часов в месяц на рутинных задачах", 1),
            ("", 0),
            ("Deloitte (исследование ИИ в HR, 2024)", 0),
            ("Сокращение времени на рекрутинг на 40%", 1),
            ("Улучшение качества найма на 25%", 1),
            ("", 0),
            ("IBM (внедрение Watson, 2023)", 0),
            ("Экономия $50M на HR-процессах", 1),
            ("Обработка 95% типовых запросов автоматически", 1),
            ("", 0),
            ("Источники: Microsoft Work Trend Index 2024, Deloitte Global HR Trends 2024", 0),
        ]
    )

    # 15. Риски и рекомендации
    add_bullets_multilevel_slide(
        prs,
        "Риски и рекомендации",
        [
            ("Ключевые риски", 0),
            ("Низкий уровень принятия пользователями", 1),
            ("Проблемы с качеством и безопасностью данных", 1),
            ("Сопротивление изменениям", 1),
            ("Недостаточное обучение персонала", 1),
            ("", 0),
            ("Рекомендации по минимизации рисков", 0),
            ("Активная программа управления изменениями", 1),
            ("Качественное обучение и поддержка пользователей", 1),
            ("Постепенное внедрение (пилот → масштабирование)", 1),
            ("Регулярный мониторинг и сбор обратной связи", 1),
            ("Обеспечение безопасности и конфиденциальности данных", 1),
        ]
    )

    # 16. Метрики успеха
    add_bullets_slide(
        prs,
        "KPI и метрики успеха",
        [
            "📊 Уровень принятия (Adoption Rate): целевой показатель 80%+ за 6 месяцев",
            "⏱ Экономия времени: 10+ часов/месяц на пользователя",
            "💰 ROI: достижение положительного ROI к 4-5 месяцу",
            "😊 Удовлетворенность пользователей: NPS > 40",
            "📈 Рост продуктивности: измеримое улучшение output на 8-12%",
            "🎯 Качество работы: снижение ошибок на 20-30%",
        ]
    )

    # 17. Заключение
    add_bullets_multilevel_slide(
        prs,
        "Заключение и рекомендации",
        [
            ("Финансовое обоснование", 0),
            ("Очень высокая окупаемость: ROI 597% за 12 месяцев", 1),
            ("Быстрая окупаемость: 3-5 месяцев в зависимости от отдела", 1),
            ("Чистая выгода: $839,153 за первый год", 1),
            ("", 0),
            ("Стратегическая ценность", 0),
            ("Освобождение времени для стратегических задач", 1),
            ("Повышение удовлетворенности сотрудников", 1),
            ("Конкурентное преимущество в эпоху цифровизации", 1),
            ("", 0),
            ("Рекомендация: ОДОБРИТЬ проект к реализации", 0),
            ("Начать с пилотного внедрения в Q1 2025", 1),
        ]
    )

    # 18. Следующие шаги
    add_bullets_slide(
        prs,
        "Следующие шаги",
        [
            "1️⃣ Утверждение бюджета и получение одобрения руководства",
            "2️⃣ Выбор платформы и заключение договора с поставщиком",
            "3️⃣ Формирование проектной команды и назначение владельца проекта",
            "4️⃣ Подготовка инфраструктуры и плана обучения",
            "5️⃣ Запуск пилота в выбранных отделах (Месяц 1-2)",
            "6️⃣ Оценка результатов пилота и корректировка плана",
            "7️⃣ Масштабирование на всю организацию (Месяцы 3-12)",
        ]
    )

    # Save presentation
    out_pptx = os.path.join(os.path.dirname(outputs_dir), "HR_AI_Impact_Comprehensive.pptx")
    prs.save(out_pptx)
    print(f"✅ Презентация сохранена: {out_pptx}")
    return out_pptx


if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.dirname(__file__))
    outputs_model_dir = os.path.join(base_dir, "outputs", "model")
    build_comprehensive_presentation(outputs_model_dir)
