#!/usr/bin/env python3
import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0


def add_image_slide(prs, title, image_path, height_inches=4.5):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(1)
    top = Inches(1.5)
    slide.shapes.add_picture(image_path, left, top, height=Inches(height_inches))


def add_table_from_csv(prs, title, csv_path, max_rows=12):
    import pandas as pd
    df = pd.read_csv(csv_path)
    if len(df) > max_rows:
        df = df.head(max_rows)
    rows, cols = df.shape
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8 + 0.3 * rows)
    table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table
    # headers
    for j, col in enumerate(df.columns):
        table.cell(0, j).text = str(col)
    # data
    for i in range(rows):
        for j in range(cols):
            table.cell(i + 1, j).text = str(df.iloc[i, j])


def build_presentation(outputs_dir: str):
    prs = Presentation()
    # Title
    add_title_slide(
        prs,
        "Автоматизация административных задач с AI-ассистентом",
        "Снижение расходов и рост продуктивности (до/после по отделам)",
    )

    # Summary bullets
    add_bullets_slide(
        prs,
        "Резюме",
        [
            "Цель: сократить затраты на административные задачи и повысить производительность",
            "Инструмент: Copilot-стиль ассистент (Generative AI/LLM)",
            "Подход: пилот → масштабирование; обучение и управление изменениями",
        ],
    )

    # Financial summary images
    add_image_slide(prs, "Кумулятивные выгоды/затраты", os.path.join(outputs_dir, "roi_over_time.png"))
    add_image_slide(prs, "Выгода по отделам", os.path.join(outputs_dir, "benefit_by_department.png"))
    add_image_slide(prs, "Затраты (12 мес)", os.path.join(outputs_dir, "costs_breakdown.png"))
    add_image_slide(prs, "Диаграмма Ганта", os.path.join(outputs_dir, "gantt.png"))

    # Department table
    add_table_from_csv(prs, "Сводка по отделам (12 мес)", os.path.join(outputs_dir, "department_summary.csv"))

    # Assumptions slide
    add_bullets_slide(
        prs,
        "Ключевые допущения",
        [
            "Лицензия: $30/польз./мес (Microsoft 365 Copilot, 2024)",
            "Доля админ. времени 20–40%; покрытие автоматизации 25–40%",
            "Монетизация продуктивности: 50% в 1-й год",
            "Обучение: 4 часа/пользователь",
        ],
    )

    # Industry references slide
    add_bullets_slide(
        prs,
        "Аналитика и примеры",
        [
            "Публикации отрасли указывают 20–30% экономии на рутинных задачах (бэк-офис)",
            "Рост throughput 5–15% при внедрении Copilot-подобных решений",
            "Сценарии: составление документов, сводки, ответы на запросы, анализ данных",
            "Риски: качество данных, безопасность, принятие пользователями",
        ],
    )

    out_pptx = os.path.join(outputs_dir, "HR_AI_Impact_Model.pptx")
    prs.save(out_pptx)
    print(f"Saved presentation: {out_pptx}")


if __name__ == "__main__":
    base_dir = os.path.dirname(os.path.dirname(__file__))
    outputs_model_dir = os.path.join(base_dir, "outputs", "model")
    build_presentation(outputs_model_dir)
